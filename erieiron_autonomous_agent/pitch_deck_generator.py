import json
import logging
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from erieiron_autonomous_agent.models import Business
from erieiron_autonomous_agent.system_agent_llm_interface import business_level_chat
from erieiron_common import common
from erieiron_common.enums import LlmModel
from erieiron_common.llm_apis.llm_interface import LlmMessage


class PitchDeckGenerator:
    """Service class to generate VC/Angel pitch decks from business data."""

    def __init__(self, business: Business):
        self.business = business

    def gather_business_data(self) -> dict:
        """Collect all relevant business data for pitch deck generation."""
        business_analysis, legal_analysis = self.business.get_latest_analysist()

        return {
            "name": self.business.name,
            "summary": self.business.summary,
            "business_plan": self.business.business_plan,
            "value_prop": self.business.value_prop,
            "revenue_model": self.business.revenue_model,
            "audience": self.business.audience,
            "core_functions": self.business.core_functions,
            "growth_channels": self.business.growth_channels,
            "ui_design_spec": self.business.ui_design_spec,
            "kpis": [
                {
                    "name": kpi.name,
                    "description": kpi.description,
                    "priority": kpi.priority
                }
                for kpi in self.business.businesskpi_set.all().order_by("-priority", "name")
            ],
            "human_work": common.get_dict(
                self.business.businesshumanjobdescription_set.all()
            ),
            "business_analysis": common.get_dict(business_analysis),
            "legal_analysis": common.get_dict(legal_analysis),
        }

    def generate_pitch_structure_with_llm(self, business_data: dict) -> dict:
        response = business_level_chat(
            description=f"Generate pitch deck structure for {self.business.name}",
            system_prompt="pitch_deck_assembler.md",
            user_messages=[
                LlmMessage.user_from_data(
                    "Business Data for Pitch Deck",
                    business_data
                )
            ],
            tag_entity=self.business,
            model=LlmModel.CLAUDE_4_5
        )

        return response

    def hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor object."""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    def create_pptx_from_structure(self, pitch_structure: dict) -> BytesIO:
        """Generate PPTX file using python-pptx with professional styling."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Extract design elements
        design = pitch_structure.get("design", {})
        primary_color = design.get("primary_color", "#1E3A8A")
        secondary_color = design.get("secondary_color", "#3B82F6")
        accent_color = design.get("accent_color", "#10B981")

        primary_rgb = self.hex_to_rgb(primary_color)
        secondary_rgb = self.hex_to_rgb(secondary_color)
        accent_rgb = self.hex_to_rgb(accent_color)

        slides_data = pitch_structure.get("slides", [])

        for idx, slide_data in enumerate(slides_data):
            title = slide_data.get("title", "")
            content = slide_data.get("content", [])

            # Use blank layout for custom styling
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add colored background bar at top
            background = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0), Inches(0),
                Inches(10), Inches(1)
            )
            background.fill.solid()
            background.fill.fore_color.rgb = primary_rgb
            background.line.fill.background()

            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.15),
                Inches(9), Inches(0.7)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            p = title_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Calibri'

            # Add slide number in corner
            slide_num_box = slide.shapes.add_textbox(
                Inches(9.2), Inches(7.1),
                Inches(0.6), Inches(0.3)
            )
            slide_num_frame = slide_num_box.text_frame
            p = slide_num_frame.paragraphs[0]
            p.text = str(idx + 1)
            p.font.size = Pt(14)
            p.font.color.rgb = primary_rgb
            p.alignment = PP_ALIGN.RIGHT

            # Add content bullets
            if content:
                content_box = slide.shapes.add_textbox(
                    Inches(0.75), Inches(1.5),
                    Inches(8.5), Inches(5.5)
                )
                text_frame = content_box.text_frame
                text_frame.word_wrap = True

                for i, bullet in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(20)
                    p.font.name = 'Calibri'
                    p.font.color.rgb = RGBColor(50, 50, 50)
                    p.space_before = Pt(12)
                    p.space_after = Pt(6)

                    # Add bullet point with accent color
                    p.font.bold = False

            # Add accent line at bottom
            accent_line = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0), Inches(7.4),
                Inches(10), Inches(0.1)
            )
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = accent_rgb
            accent_line.line.fill.background()

        # Save to BytesIO buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        return buffer

    def generate_pitch_deck(self) -> BytesIO:
        """Main orchestration method to generate complete pitch deck."""
        logging.info(f"Generating pitch deck for business: {self.business.name}")

        # Gather business data
        business_data = self.gather_business_data()
        logging.info(f"Gathered business data with {len(business_data)} fields")

        # Generate structure via LLM
        pitch_structure = self.generate_pitch_structure_with_llm(business_data)
        slide_count = len(pitch_structure.get("slides", []))
        logging.info(f"Generated pitch deck structure with {slide_count} slides")

        # Create PPTX file
        pptx_buffer = self.create_pptx_from_structure(pitch_structure)
        logging.info(f"Created PPTX file, size: {len(pptx_buffer.getvalue())} bytes")

        return pptx_buffer
